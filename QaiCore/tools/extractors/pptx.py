"""
Extractor de texto de archivos PPTX (PowerPoint).

Dependencias: pip install python-pptx
"""

from pptx import Presentation
from pathlib import Path


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extrae texto de un archivo PPTX.
    
    Args:
        pptx_path: Ruta al archivo PPTX
        
    Returns:
        Texto extraído de todas las diapositivas
    """
    try:
        prs = Presentation(pptx_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"\n{'='*60}", f"SLIDE {slide_num}", f"{'='*60}"]
            
            # Extraer texto de todas las formas de la diapositiva
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extraer texto de tablas si existen
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_text.append(row_text)
            
            if len(slide_text) > 3:  # Solo agregar si hay contenido más allá del header
                text_parts.extend(slide_text)
        
        return "\n".join(text_parts)
        
    except Exception as e:
        return f"❌ Error extrayendo texto del PPTX: {str(e)}"


# Ejemplo de uso
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Uso: python pptx.py <ruta_archivo.pptx>")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    
    if not Path(pptx_file).exists():
        print(f"❌ Archivo no encontrado: {pptx_file}")
        sys.exit(1)
    
    print(f"📊 Extrayendo texto de: {pptx_file}")
    text = extract_text_from_pptx(pptx_file)
    
    print("\n" + text[:800] if len(text) > 800 else text)
    print("\n" + "="*60)
    print(f"✅ Total extraído: {len(text)} caracteres")
